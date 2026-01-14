# File parsing utilities for supported formats
# Supported file types:
# - .txt (plain text)
# - .csv (comma-separated values)
# - .xls, .xlsx (Excel spreadsheets)
# - .docx (Word documents)
# - .pptx (PowerPoint presentations, modern format)
# - .ppt (PowerPoint presentations, legacy format - requires pywin32 on Windows)
# - .pdf (Portable Document Format with OCR support for images)

import os
import pandas as pd
import docx
import pptx
import tempfile
import io
from supabase import create_client
import base64
from datetime import datetime
from PIL import Image
import io as io_module

# Supabase configuration
SUPABASE_URL = os.environ.get("NEXT_PUBLIC_SUPABASE_URL", "https://fmlanqjduftxlktygpwe.supabase.co")
SUPABASE_KEY = os.environ.get("SUPABASE_SERVICE_ROLE_KEY", "eyJhbGciOiJIUzI1NiIsInR5cCI6IkpXVCJ9.eyJpc3MiOiJzdXBhYmFzZSIsInJlZiI6ImZtbGFucWpkdWZ0eGxrdHlncHdlIiwicm9sZSI6ImFub24iLCJpYXQiOjE3NTk0MDkzNTcsImV4cCI6MjA3NDk4NTM1N30.FT6c6BNfkJJFKliI1qv9uzBJj0UWMIaykRJrwKQKIfs")

# Use pypdf instead of deprecated PyPDF2
try:
    from pypdf import PdfReader
    PDF_LIBRARY = "pypdf"
except ImportError:
    try:
        import PyPDF2
        PDF_LIBRARY = "PyPDF2"
        print("Warning: Using deprecated PyPDF2. Consider upgrading to pypdf: pip install pypdf")
    except ImportError:
        PDF_LIBRARY = None
        print("Warning: No PDF library available. Install pypdf: pip install pypdf")

# OCR support for image extraction
try:
    from pdf2image import convert_from_path
    PDF2IMAGE_AVAILABLE = True
except ImportError:
    PDF2IMAGE_AVAILABLE = False
    print("Warning: pdf2image not available. Install it for better PDF image extraction: pip install pdf2image")

try:
    import pytesseract
    PYTESSERACT_AVAILABLE = True
except ImportError:
    PYTESSERACT_AVAILABLE = False
    print("Warning: pytesseract not available. Install it for OCR support: pip install pytesseract")

def extract_text_from_image(image_path: str):
    """
    Extract text from an image using OCR (Tesseract)
    Args:
        image_path: Path to the image file
    Returns:
        Extracted text from the image
    """
    if not PYTESSERACT_AVAILABLE:
        return ""
    
    try:
        image = Image.open(image_path)
        text = pytesseract.image_to_string(image)
        return text.strip()
    except Exception as e:
        print(f"Warning: Failed to extract text from image {image_path}: {e}")
        return ""

def extract_text_from_pdf_with_ocr(file_path: str):
    """
    Extract text from PDF, including text from images using OCR
    Args:
        file_path: Path to the PDF file
    Returns:
        Combined extracted text from all pages
    """
    text_parts = []
    
    # First, try to extract text using pypdf (faster for text-based PDFs)
    try:
        if PDF_LIBRARY == "pypdf":
            with open(file_path, "rb") as f:
                reader = PdfReader(f)
                for page_num, page in enumerate(reader.pages, 1):
                    extracted_text = page.extract_text()
                    if extracted_text and extracted_text.strip():
                        text_parts.append(f"--- Page {page_num} (Text) ---")
                        text_parts.append(extracted_text)
        elif PDF_LIBRARY == "PyPDF2":
            with open(file_path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                for page_num, page in enumerate(reader.pages, 1):
                    extracted_text = page.extract_text()
                    if extracted_text and extracted_text.strip():
                        text_parts.append(f"--- Page {page_num} (Text) ---")
                        text_parts.append(extracted_text)
    except Exception as e:
        print(f"Warning: Failed to extract text from PDF with pypdf: {e}")
    
    # Then, use OCR on images for pages with little/no text or to capture image content
    if PDF2IMAGE_AVAILABLE and PYTESSERACT_AVAILABLE:
        try:
            print("🔍 Attempting OCR extraction from PDF images...")
            images = convert_from_path(file_path)
            
            for page_num, image in enumerate(images, 1):
                # Check if we already extracted good text from this page
                existing_text = ""
                for part in text_parts:
                    if f"Page {page_num}" in part:
                        existing_text = part
                        break
                
                # Only do OCR if we didn't get much text already
                if not existing_text or len(existing_text) < 100:
                    try:
                        ocr_text = pytesseract.image_to_string(image)
                        if ocr_text and ocr_text.strip():
                            text_parts.append(f"--- Page {page_num} (OCR) ---")
                            text_parts.append(ocr_text)
                            print(f"✅ Extracted {len(ocr_text)} characters from page {page_num} via OCR")
                    except Exception as e:
                        print(f"Warning: OCR failed for page {page_num}: {e}")
        except Exception as e:
            print(f"Warning: PDF to image conversion failed: {e}")
    
    return "\n".join(text_parts) if text_parts else ""

def extract_text_from_docx_with_images(file_path: str):
    """
    Extract text from DOCX including text from embedded images
    Args:
        file_path: Path to the DOCX file
    Returns:
        Combined extracted text and OCR text from images
    """
    text_parts = []
    
    try:
        doc = docx.Document(file_path)
        
        # Extract text from paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        
        # Extract images and perform OCR
        if PYTESSERACT_AVAILABLE:
            # Extract images from the document
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    try:
                        image_part = rel.target_part
                        image_bytes = image_part.blob
                        
                        # Create temporary image file
                        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp_img:
                            tmp_img.write(image_bytes)
                            tmp_img_path = tmp_img.name
                        
                        # Extract text from image
                        ocr_text = extract_text_from_image(tmp_img_path)
                        if ocr_text:
                            text_parts.append(f"\n[Image Content]\n{ocr_text}")
                        
                        # Clean up
                        os.remove(tmp_img_path)
                    except Exception as e:
                        print(f"Warning: Failed to process image in DOCX: {e}")
        
        return "\n".join(text_parts)
    except Exception as e:
        print(f"Error extracting text from DOCX: {str(e)}")
        raise

def store_extracted_content(file_id: str, user_id: str, content: str, file_name: str, file_path: str = None):
    """
    Store extracted text content in the document_content table
    Args:
        file_id: ID of the file in file_upload table
        user_id: User ID who uploaded the file
        content: Extracted text content
        file_name: Original file name
        file_path: Path to the original file in Supabase storage or local filesystem
    Returns True on success, False on failure
    """
    try:
        supabase = create_client(SUPABASE_URL, SUPABASE_KEY)
        
        # Upsert the content (insert or update if exists)
        result = supabase.table('document_content').upsert({
            'file_id': file_id,
            'user_id': user_id,
            'content': content.strip(),
            'file_name': file_name,
            'file_path': file_path,  # Store reference to original file
            'extracted_at': datetime.utcnow().isoformat()
        }, on_conflict='file_id').execute()
        
        # Check if upsert was successful
        if result and result.data:
            print(f"✅ Successfully stored extracted content for file: {file_name}")
            print(f"   Upserted record: {result.data}")
            return True
        else:
            print(f"⚠️  Warning: Upsert returned no data. Response: {result}")
            return False
            
    except Exception as e:
        print(f"❌ Error storing extracted content: {str(e)}")
        return False


def extract_and_store_file_content(file_path: str, file_id: str, user_id: str, file_name: str) -> tuple[str, bool]:
    """
    Extract text from file and store it in Supabase
    Returns (extracted_text, storage_success)
    """
    try:
        extracted_text = extract_text_from_file(file_path)
        # Pass file_path as reference for later structured data access
        storage_success = store_extracted_content(file_id, user_id, extracted_text, file_name, file_path=file_path)
        return extracted_text, storage_success
    except Exception as e:
        print(f"Error in extract_and_store_file_content: {str(e)}")
        return "", False

def extract_text_from_file(file_path: str):
    """Extract text from various file formats"""
    print(f"Attempting to parse file: {file_path}")
    downloaded_temp_path = None

    # If the file doesn't exist locally, try to download it from Supabase Storage
    if not os.path.exists(file_path):
        if SUPABASE_URL and SUPABASE_KEY:
            try:
                print("File not found locally. Attempting to download from Supabase storage...")
                supabase = create_client(SUPABASE_URL, SUPABASE_KEY)
                # Assume bucket name is 'vault_files' and file_path is the object key
                file_data = supabase.storage.from_('vault_files').download(file_path)

                # The python-supabase client may return bytes or a response-like object
                if hasattr(file_data, 'read'):
                    file_bytes = file_data.read()
                else:
                    file_bytes = file_data

                if not file_bytes:
                    raise ValueError("Downloaded file is empty or could not be retrieved from Supabase")

                # Write to a temporary file so existing parsing logic can operate on a path
                with tempfile.NamedTemporaryFile(delete=False, suffix="_" + os.path.basename(file_path)) as tf:
                    tf.write(file_bytes)
                    downloaded_temp_path = tf.name

                print(f"Downloaded file to temporary path: {downloaded_temp_path}")
                # Use the temporary file for subsequent parsing
                file_path = downloaded_temp_path
            except Exception as e:
                print(f"Warning: Failed to download file from Supabase: {e}")
                # Fall through and let later logic raise a helpful error when file is missing
        else:
            print("Supabase credentials (SUPABASE_URL / SUPABASE_SERVICE_ROLE_KEY) not set; cannot download missing file")
    
    # Get file basename for dotfile checking
    base_name = os.path.basename(file_path).lower()
    
    # Get file extension (handle cases like .pptx.pptx)
    ext = os.path.splitext(file_path)[1].lower()
    
    # If file has double extension, use the real extension
    if ext == '.pptx' or file_path.lower().endswith('.pptx.pptx'):
        ext = '.pptx'
    elif ext == '.ppt' or file_path.lower().endswith('.ppt.ppt'):
        ext = '.ppt'
    
    print(f"Detected file extension: {ext}")
    
    try:
        if ext == ".csv":
            # Load CSV and extract meaningful metadata for semantic search
            try:
                df = pd.read_csv(file_path)
            except UnicodeDecodeError:
                df = pd.read_csv(file_path, encoding='latin-1')
            except Exception as e:
                print(f"Error reading CSV with pandas: {str(e)}")
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()
            
            # Build meaningful content for semantic search
            text_parts = []
            
            # Add file metadata and column information
            text_parts.append(f"CSV File: {os.path.basename(file_path)}")
            text_parts.append(f"Rows: {len(df)}, Columns: {len(df.columns)}")
            text_parts.append(f"Column Names: {', '.join(df.columns)}")
            text_parts.append("-" * 80)
            
            # Add column data types
            text_parts.append("Column Types:")
            for col in df.columns:
                text_parts.append(f"  {col}: {df[col].dtype}")
            text_parts.append("-" * 80)
            
            # Add first 5 rows as sample (helps semantic search identify content)
            text_parts.append("Data Sample (first 5 rows):")
            text_parts.append(df.head(5).to_string())
            
            return "\n".join(text_parts)
        
        elif ext in [".xls", ".xlsx"]:
            # Load Excel and extract meaningful metadata for semantic search
            try:
                xls = pd.ExcelFile(file_path)
                sheet_names = xls.sheet_names
                text_parts = []
                
                text_parts.append(f"Excel File: {os.path.basename(file_path)}")
                text_parts.append(f"Total Sheets: {len(sheet_names)}")
                text_parts.append(f"Sheet Names: {', '.join(sheet_names)}")
                text_parts.append("=" * 80)
                
                # Process first sheet only for metadata (to keep embedding size reasonable)
                if sheet_names:
                    try:
                        df = pd.read_excel(file_path, sheet_name=sheet_names[0])
                        
                        text_parts.append(f"\nPrimary Sheet: {sheet_names[0]}")
                        text_parts.append(f"Rows: {len(df)}, Columns: {len(df.columns)}")
                        text_parts.append(f"Column Names: {', '.join(df.columns)}")
                        text_parts.append("-" * 80)
                        
                        text_parts.append("Column Types:")
                        for col in df.columns:
                            text_parts.append(f"  {col}: {df[col].dtype}")
                        text_parts.append("-" * 80)
                        
                        # Add first 5 rows as sample
                        text_parts.append("Data Sample (first 5 rows):")
                        text_parts.append(df.head(5).to_string())
                    except Exception as e:
                        text_parts.append(f"Error reading primary sheet: {str(e)}")
                
                return "\n".join(text_parts)
            
            except Exception as e:
                print(f"Error reading Excel file: {str(e)}")
                try:
                    df = pd.read_excel(file_path, sheet_name=0)
                    return df.to_string()
                except Exception as fallback_error:
                    raise ValueError(f"Unable to parse Excel file: {str(fallback_error)}")
        elif ext == ".docx":
            print("📄 Extracting text from DOCX with image OCR support...")
            return extract_text_from_docx_with_images(file_path)
        elif ext == ".pptx":
            print("🖼️  Extracting text from PPTX with image OCR support...")
            prs = pptx.Presentation(file_path)
            text = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                text.append(f"\n--- Slide {slide_num} ---")
                
                # Extract text from all shapes using improved method
                for shape in slide.shapes:
                    shape_texts = extract_text_from_shape(shape)
                    text.extend(shape_texts)
                    
                    # Extract images from shapes and perform OCR
                    if PYTESSERACT_AVAILABLE and hasattr(shape, "image"):
                        try:
                            image = shape.image
                            image_bytes = image.blob
                            
                            # Create temporary image file
                            with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp_img:
                                tmp_img.write(image_bytes)
                                tmp_img_path = tmp_img.name
                            
                            # Extract text from image
                            ocr_text = extract_text_from_image(tmp_img_path)
                            if ocr_text:
                                text.append(f"[Image OCR] {ocr_text}")
                            
                            # Clean up
                            os.remove(tmp_img_path)
                        except Exception as e:
                            print(f"Warning: Failed to extract image from slide {slide_num}: {e}")
                
                # Extract notes
                try:
                    if slide.has_notes_slide:
                        notes_text = slide.notes_slide.notes_text_frame.text
                        if notes_text and notes_text.strip():
                            text.append(f"[Notes] {notes_text.strip()}")
                except Exception as e:
                    print(f"Warning: Error extracting notes from slide: {e}")
            
            content = "\n".join(text)
            print(f"✅ Extracted {len(content)} characters from PowerPoint")
            return content
        elif ext == ".ppt":
            # Handle older .ppt format (binary PowerPoint files)
            # Try multiple approaches for better compatibility
            
            # Approach 1: Try LibreOffice conversion (if available)
            try:
                import subprocess
                import tempfile
                
                # Create a temporary directory for conversion
                with tempfile.TemporaryDirectory() as temp_dir:
                    # Try to convert using LibreOffice
                    output_file = os.path.join(temp_dir, "converted.pptx")
                    
                    # Try common LibreOffice paths
                    libreoffice_paths = [
                        r"C:\Program Files\LibreOffice\program\soffice.exe",
                        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
                        "soffice",  # If in PATH
                        "libreoffice",  # Alternative command
                    ]
                    
                    conversion_successful = False
                    for lo_path in libreoffice_paths:
                        try:
                            result = subprocess.run(
                                [lo_path, "--headless", "--convert-to", "pptx", 
                                 "--outdir", temp_dir, file_path],
                                capture_output=True,
                                text=True,
                                timeout=30
                            )
                            
                            # Check if conversion created a .pptx file
                            converted_files = [f for f in os.listdir(temp_dir) if f.endswith('.pptx')]
                            if converted_files:
                                converted_path = os.path.join(temp_dir, converted_files[0])
                                # Now parse the converted .pptx file
                                prs = pptx.Presentation(converted_path)
                                text = []
                                
                                for slide in prs.slides:
                                    for shape in slide.shapes:
                                        shape_texts = extract_text_from_shape(shape)
                                        text.extend(shape_texts)
                                
                                content = "\n".join(text)
                                print(f"Extracted {len(content)} characters from .ppt file (via LibreOffice conversion)")
                                return content
                        except (subprocess.TimeoutExpired, subprocess.CalledProcessError, FileNotFoundError):
                            continue
                    
                    # If we get here, LibreOffice conversion didn't work
                    print("LibreOffice conversion not available or failed")
                    
            except Exception as e:
                print(f"LibreOffice approach failed: {str(e)}")
            
            # Approach 2: Try pywin32 COM automation (Windows only)
            try:
                import win32com.client
                import pythoncom
                
                pythoncom.CoInitialize()
                try:
                    powerpoint = win32com.client.Dispatch("PowerPoint.Application")
                    # Don't set Visible property - let PowerPoint decide
                    # Some versions don't allow hiding the application
                    try:
                        powerpoint.DisplayAlerts = 0  # Use 0 instead of False for COM
                    except:
                        pass  # Some versions don't support this
                    
                    # Convert to absolute path with proper formatting
                    abs_path = os.path.abspath(file_path)
                    
                    # Open with more permissive settings
                    # Use integer constants instead of boolean for COM compatibility
                    presentation = powerpoint.Presentations.Open(
                        abs_path, 
                        ReadOnly=1,  # Use 1 for True in COM
                        Untitled=0,  # Use 0 for False in COM
                        WithWindow=0  # Use 0 for False in COM
                    )
                    
                    text = []
                    skipped_shapes = 0
                    try:
                        for slide_idx in range(1, presentation.Slides.Count + 1):
                            slide = presentation.Slides(slide_idx)
                            for shape_idx in range(1, slide.Shapes.Count + 1):
                                try:
                                    shape = slide.Shapes(shape_idx)
                                    if hasattr(shape, "TextFrame") and hasattr(shape.TextFrame, "TextRange"):
                                        shape_text = shape.TextFrame.TextRange.Text
                                        if shape_text and shape_text.strip():
                                            text.append(shape_text.strip())
                                except Exception as shape_error:
                                    # Common errors for non-text shapes (images, charts, grouped objects)
                                    # These are expected and can be safely ignored
                                    skipped_shapes += 1
                                    continue
                    finally:
                        presentation.Close()
                        powerpoint.Quit()
                    
                    content = "\n".join(text)
                    if skipped_shapes > 0:
                        print(f"Extracted {len(content)} characters from .ppt file (via COM). Skipped {skipped_shapes} non-text shapes (images/charts/groups).")
                    else:
                        print(f"Extracted {len(content)} characters from .ppt file (via COM)")
                    return content
                    
                finally:
                    pythoncom.CoUninitialize()
                
            except ImportError:
                raise ValueError(
                    ".ppt files require either:\n"
                    "1. LibreOffice installed (recommended, cross-platform)\n"
                    "   Download from: https://www.libreoffice.org/download/\n"
                    "2. pywin32 on Windows: pip install pywin32\n"
                    "3. Or convert .ppt to .pptx format manually using PowerPoint"
                )
            except Exception as e:
                raise ValueError(
                    f"Error reading .ppt file: {str(e)}\n\n"
                    "Solutions:\n"
                    "1. Install LibreOffice from https://www.libreoffice.org/download/\n"
                    "2. Convert the file to .pptx format using PowerPoint or online converter\n"
                    "3. Ensure PowerPoint is properly installed (for Windows COM approach)"
                )
        elif ext == ".pdf":
            print("📄 Extracting text from PDF with OCR support for images...")
            return extract_text_from_pdf_with_ocr(file_path)
        elif ext == ".txt":
            with open(file_path, "r", encoding="utf-8") as f:
                return f.read()
        elif ext in [".py", ".js", ".ts", ".java", ".cpp", ".c", ".h", ".cs", ".go", ".rs", 
                     ".html", ".css", ".scss", ".jsx", ".tsx", ".json", ".yaml", ".yml", 
                     ".toml", ".ini", ".env", ".md", ".sh", ".bat", ".ps1", ".sql", ".prisma", ".graphql"] or \
             base_name in [".dockerignore", ".gitignore"]:
            # Source code, configuration files, and dotfiles - read as text
            try:
                with open(file_path, "r", encoding="utf-8") as f:
                    return f.read()
            except UnicodeDecodeError:
                # Fallback to latin-1 for files with encoding issues
                with open(file_path, "r", encoding="latin-1") as f:
                    return f.read()
        else:
            raise ValueError(f"Unsupported file type: {ext}")
    except Exception as e:
        print(f"Error parsing file {file_path}: {str(e)}")
        raise
    finally:
        # Clean up any temporary file we downloaded from Supabase
        try:
            if downloaded_temp_path and os.path.exists(downloaded_temp_path):
                os.remove(downloaded_temp_path)
                print(f"Cleaned up temporary file: {downloaded_temp_path}")
        except Exception:
            pass

def extract_text_from_shape(shape):
    """Recursively extract text from a shape and its sub-shapes (for PPTX files)"""
    shape_texts = []
    
    try:
        # Method 1: Direct text access
        if hasattr(shape, 'text') and shape.text and shape.text.strip():
            shape_texts.append(shape.text.strip())
        
        # Method 2: Text frame access with paragraphs and runs
        if hasattr(shape, 'text_frame') and shape.text_frame:
            if hasattr(shape.text_frame, 'paragraphs'):
                for para in shape.text_frame.paragraphs:
                    if hasattr(para, 'text') and para.text and para.text.strip():
                        para_text = para.text.strip()
                        if para_text not in shape_texts:
                            shape_texts.append(para_text)
                    
                    # Extract from runs within paragraphs
                    if hasattr(para, 'runs'):
                        for run in para.runs:
                            if hasattr(run, 'text') and run.text and run.text.strip():
                                run_text = run.text.strip()
                                if run_text not in shape_texts:
                                    shape_texts.append(run_text)
        
        # Method 3: Handle GROUP shapes recursively
        if hasattr(shape, 'shapes'):  # This is a group
            for sub_shape in shape.shapes:
                sub_texts = extract_text_from_shape(sub_shape)
                shape_texts.extend(sub_texts)
        
        # Method 4: Handle tables
        if hasattr(shape, 'has_table') and shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text and cell.text.strip():
                        cell_text = cell.text.strip()
                        if cell_text not in shape_texts:
                            shape_texts.append(cell_text)
    
    except Exception as e:
        print(f"Warning: Error extracting text from shape: {e}")
    
    return shape_texts

"""
Google Drive API Async Wrapper

Provides async functions for searching Drive files and reading document content.
Uses Google Drive REST API v3.
"""

import httpx
from typing import List, Dict, Any, Optional


# MIME type mapping for Google Workspace files
GOOGLE_WORKSPACE_MIME_TYPES = {
    "document": "application/vnd.google-apps.document",
    "spreadsheet": "application/vnd.google-apps.spreadsheet",
    "presentation": "application/vnd.google-apps.presentation",
    "pdf": "application/pdf",
}

# Export formats for Google Workspace files
EXPORT_MIME_TYPES = {
    "application/vnd.google-apps.document": "text/plain",
    "application/vnd.google-apps.spreadsheet": "text/csv",
    "application/vnd.google-apps.presentation": "text/plain",
}

BASE_URL = "https://www.googleapis.com/drive/v3"
TIMEOUT = 30.0


async def search_drive(
    query: str,
    access_token: str,
    max_results: int = 10,
    file_type: Optional[str] = None,
):
    """Search Google Drive for files matching a query.

    Args:
        query: Search query (Google Drive query syntax or natural text)
        access_token: OAuth2 access token
        max_results: Maximum number of results to return
        file_type: Optional filter (document, spreadsheet, presentation, pdf)

    Returns:
        List of file result dicts with title, snippet, url, date, file_type
    """
    headers = {"Authorization": f"Bearer {access_token}"}

    # Build Drive API query string
    drive_query_parts = []

    # If query looks like Drive query syntax, validate and use it directly
    if any(op in query for op in ["name contains", "fullText contains", "mimeType ="]):
        # Query appears to be properly formatted Drive API syntax - validate it
        import re
        
        # Check for obviously malformed syntax that should be rejected
        has_invalid_syntax = (
            "~" in query or  # tilde is not a valid operator
            ", " in query and " and " not in query.lower() or  # comma-separated without 'and'
            re.search(r"mimeType\s*=\s*'[^']*[^/]'", query)  # MIME type with trailing non-slash
        )
        
        if not has_invalid_syntax:
            normalized_query = query.replace('"', "'")  # Normalize quotes
            
            # Check if query already contains 'trashed = false'
            if "trashed = false" not in normalized_query:
                drive_query_parts.append(normalized_query)
            else:
                drive_query = normalized_query
                drive_query_parts = []
        else:
            # Query has invalid syntax, treat as natural language search instead
            escaped_query = _escape_query(query)
            if " " in query.strip():
                drive_query_parts.append(f"fullText contains '{escaped_query}'")
            else:
                drive_query_parts.append(f"fullText contains '{escaped_query}'")
    elif query.strip():  # Only add fullText search if query is not empty
        escaped_query = _escape_query(query)
        # Use fullText contains for all natural language queries (most reliable)
        drive_query_parts.append(f"fullText contains '{escaped_query}'")

    # Apply file type filter
    if file_type and file_type in GOOGLE_WORKSPACE_MIME_TYPES:
        drive_query_parts.append(f"mimeType = '{GOOGLE_WORKSPACE_MIME_TYPES[file_type]}'")

    # Exclude trashed files (only if not already set above)
    if not (any(op in query for op in ["name contains", "fullText contains", "mimeType ="]) and "trashed = false" in query):
        drive_query_parts.append("trashed = false")

    # Join query parts
    if drive_query_parts:
        drive_query = " and ".join(drive_query_parts)
    # If drive_query was set above (for queries with existing trashed filter), keep it

    params = {
        "q": drive_query,
        "pageSize": min(max_results, 50),
        "fields": "files(id,name,mimeType,modifiedTime,webViewLink,description,size,owners)",
    }
    
    # Use modifiedTime desc as default - it's more reliable than relevance
    # Only use relevance for simple single-word fullText searches
    if "fullText contains" in drive_query and len(drive_query.split("fullText contains")) == 2:
        # Simple single fullText query - try relevance, fallback to modifiedTime if needed
        search_term = drive_query.split("fullText contains")[1].strip()
        if search_term.count("'") == 2 and " " not in search_term.replace("'", "").strip():
            params["orderBy"] = "relevance"
        else:
            params["orderBy"] = "modifiedTime desc"
    else:
        params["orderBy"] = "modifiedTime desc"

    async with httpx.AsyncClient(timeout=TIMEOUT) as client:
        resp = await client.get(
            f"{BASE_URL}/files",
            headers=headers,
            params=params,
        )
        resp.raise_for_status()
        data = resp.json()

    files = data.get("files", [])
    results = []

    for f in files:
        owner = ""
        if f.get("owners"):
            owner = f["owners"][0].get("displayName", "")

        result = {
            "title": f.get("name", "Untitled"),
            "url": f.get("webViewLink", ""),
            "date": f.get("modifiedTime", ""),
            "content": f.get("description", ""),
            "file_type": _friendly_mime_type(f.get("mimeType", "")),
            "mime_type": f.get("mimeType", ""),  # Store actual MIME type for content extraction
            "author": owner,
            "file_id": f.get("id"),
            "size": f.get("size"),
        }
        results.append(result)

    # Optionally fetch content snippets for top results
    if results and len(results) <= 5:
        for result in results[:3]:  # Read content of top 3 files
            try:
                print(f"📄 Fetching content for: {result['title']} (mime_type={result.get('mime_type')})")
                content = await get_file_content(
                    result["file_id"], access_token, result.get("mime_type")
                )
                if content:
                    result["content"] = content[:2000]  # Truncate for context
                    print(f"✅ Content extracted ({len(content)} chars): {content[:100]}...")
                else:
                    print(f"⚠️  No content extracted for {result['title']}")
            except Exception as e:
                print(f"❌ Content fetch error for {result['title']}: {str(e)}")
                pass  # Content fetch is best-effort

    return results


async def get_file_content(
    file_id: str, access_token: str, mime_type: Optional[str] = None
):
    """Read the text content of a Google Drive file.

    Supports multiple file formats:
    - Google Workspace files (Docs, Sheets, Slides) → exported as text
    - PDFs → text extraction with pdfplumber
    - Word (.docx) → text extraction with python-docx
    - Excel (.xlsx) → text extraction with openpyxl
    - PowerPoint (.pptx) → text extraction with python-pptx
    - Plain text, CSV, etc. → direct text reading

    Args:
        file_id: The Google Drive file ID
        access_token: OAuth2 access token
        mime_type: Optional MIME type hint

    Returns:
        File content as text, or None on failure
    """
    headers = {"Authorization": f"Bearer {access_token}"}

    async with httpx.AsyncClient(timeout=TIMEOUT) as client:
        # First, get file metadata to determine type
        if not mime_type:
            try:
                meta_resp = await client.get(
                    f"{BASE_URL}/files/{file_id}",
                    headers=headers,
                    params={"fields": "mimeType"},
                )
                meta_resp.raise_for_status()
                mime_type = meta_resp.json().get("mimeType", "")
            except Exception as e:
                print(f"❌ Failed to get file metadata: {str(e)}")
                return None

        # Google Workspace files need to be exported
        if mime_type in EXPORT_MIME_TYPES:
            export_mime = EXPORT_MIME_TYPES[mime_type]
            try:
                resp = await client.get(
                    f"{BASE_URL}/files/{file_id}/export",
                    headers=headers,
                    params={"mimeType": export_mime},
                )
                resp.raise_for_status()
                if resp.status_code == 200:
                    return resp.text[:50000]
            except Exception as e:
                print(f"❌ Failed to export {mime_type}: {str(e)}")
                return None
        else:
            # Regular files - download content
            try:
                resp = await client.get(
                    f"{BASE_URL}/files/{file_id}",
                    headers=headers,
                    params={"alt": "media"},
                )
                resp.raise_for_status()
                
                if resp.status_code == 200:
                    return _extract_text_from_bytes(resp.content, mime_type)
            except Exception as e:
                print(f"❌ Failed to download file: {str(e)}")
                return None

        return None


def _extract_text_from_bytes(content: bytes, mime_type: str) -> Optional[str]:
    """Extract text from various file formats."""
    import io
    
    try:
        # PDF files
        if mime_type and "pdf" in mime_type.lower():
            try:
                import pdfplumber
                pdf_bytes = io.BytesIO(content)
                with pdfplumber.open(pdf_bytes) as pdf:
                    text = ""
                    # Extract text from all pages (up to limit)
                    for page_num, page in enumerate(pdf.pages):
                        text += page.extract_text() or ""
                        if len(text) > 50000:
                            break
                print(f"✅ PDF: Extracted {len(text)} chars from {len(pdf.pages)} pages")
                return text[:50000]
            except ImportError:
                print(f"⚠️  pdfplumber not installed for PDF extraction")
                return f"[PDF Document - {len(content)} bytes]"
            except Exception as e:
                print(f"⚠️  Failed to extract PDF text: {str(e)}")
                return None
        
        # Word files (.docx)
        elif mime_type and ("word" in mime_type.lower() or "docx" in mime_type.lower() or "wordprocessingml" in mime_type.lower()):
            try:
                from docx import Document
                doc_bytes = io.BytesIO(content)
                doc = Document(doc_bytes)
                text = "\n".join([para.text for para in doc.paragraphs])
                print(f"✅ DOCX: Extracted {len(text)} chars from {len(doc.paragraphs)} paragraphs")
                return text[:50000]
            except ImportError:
                print(f"⚠️  python-docx not installed for DOCX extraction")
                return None
            except Exception as e:
                print(f"⚠️  Failed to extract DOCX text: {str(e)}")
                return None
        
        # Excel files (.xlsx)
        elif mime_type and ("excel" in mime_type.lower() or "spreadsheet" in mime_type.lower() or "sheet" in mime_type.lower()):
            try:
                from openpyxl import load_workbook
                excel_bytes = io.BytesIO(content)
                wb = load_workbook(excel_bytes)
                text = ""
                # Extract from all sheets (up to limit)
                for sheet_idx, sheet in enumerate(wb.sheetnames):
                    ws = wb[sheet]
                    text += f"\n=== Sheet: {sheet} ===\n"
                    for row in ws.iter_rows(values_only=True):
                        text += " | ".join(str(cell) if cell else "" for cell in row) + "\n"
                        if len(text) > 50000:
                            break
                    if len(text) > 50000:
                        break
                print(f"✅ XLSX: Extracted {len(text)} chars from {len(wb.sheetnames)} sheets")
                return text[:50000]
            except ImportError:
                print(f"⚠️  openpyxl not installed for XLSX extraction")
                return None
            except Exception as e:
                print(f"⚠️  Failed to extract XLSX text: {str(e)}")
                return None
        
        # PowerPoint files (.pptx)
        elif mime_type and ("powerpoint" in mime_type.lower() or "presentation" in mime_type.lower() or "pptx" in mime_type.lower()):
            try:
                from pptx import Presentation
                ppt_bytes = io.BytesIO(content)
                prs = Presentation(ppt_bytes)
                text = ""
                # Extract text from all slides (up to limit)
                for slide_idx, slide in enumerate(prs.slides):
                    text += f"\n=== Slide {slide_idx + 1} ===\n"
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                    if len(text) > 50000:
                        break
                print(f"✅ PPTX: Extracted {len(text)} chars from {len(prs.slides)} slides")
                return text[:50000]
            except ImportError:
                print(f"⚠️  python-pptx not installed for PPTX extraction")
                return None
            except Exception as e:
                print(f"⚠️  Failed to extract PPTX text: {str(e)}")
                return None
        
        # Text files, CSV, JSON, etc.
        else:
            try:
                text = content.decode('utf-8', errors='ignore')
                print(f"✅ Text file: Extracted {len(text)} chars")
                return text[:50000]
            except Exception as e:
                print(f"⚠️  Failed to decode as text: {str(e)}")
                return f"[Binary file - {len(content)} bytes]"
    
    except Exception as e:
        print(f"❌ Unexpected error in text extraction: {str(e)}")
        return None


def _escape_query(query: str):
    """Escape single quotes for Google Drive query syntax."""
    return query.replace("'", "\\'")


def _friendly_mime_type(mime_type: str):
    """Convert MIME type to a user-friendly label."""
    mapping = {
        "application/vnd.google-apps.document": "Google Doc",
        "application/vnd.google-apps.spreadsheet": "Google Sheet",
        "application/vnd.google-apps.presentation": "Google Slides",
        "application/pdf": "PDF",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "Word",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "Excel",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": "PowerPoint",
        "text/plain": "Text",
        "image/png": "Image (PNG)",
        "image/jpeg": "Image (JPEG)",
    }
    return mapping.get(mime_type, mime_type.split("/")[-1] if "/" in mime_type else mime_type)

"""
Presentation Generator Agent
Generates professional presentations on any topic using python-pptx and LLM-generated content
"""

import os
import json
import requests
from datetime import datetime
from typing import Optional, Dict, Any, List
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from dotenv import load_dotenv

# Load environment variables
env_path = os.path.join(os.path.dirname(__file__), '.env.local')
load_dotenv(dotenv_path=env_path)

# LLM Configuration
OLLAMA_API_URL = os.getenv("OLLAMA_API_URL", "http://localhost:11434")
OLLAMA_MODEL = os.getenv("OLLAMA_MODEL", "llama3.2:3b")

# Color scheme for presentations
COLOR_PRIMARY = RGBColor(31, 78, 121)  # Dark Blue
COLOR_SECONDARY = RGBColor(79, 129, 189)  # Light Blue
COLOR_ACCENT = RGBColor(192, 0, 0)  # Red
COLOR_TEXT = RGBColor(51, 51, 51)  # Dark Gray
COLOR_WHITE = RGBColor(255, 255, 255)  # White


def query_llm_for_presentation(topic: str, num_slides: int = 10, style: str = "professional") -> Dict[str, Any]:
    """
    Use LLM to generate presentation outline and content
    
    Args:
        topic: Topic for the presentation
        num_slides: Number of slides to generate (default: 10)
        style: Presentation style (professional, educational, creative)
    
    Returns:
        Dictionary with presentation structure including title, slides, and speaker notes
    """
    try:
        prompt = f"""Generate a professional presentation outline on the topic: "{topic}"

Create exactly {num_slides} slides with the following structure for each slide:
1. Title slide: Include a catchy title and subtitle
2-{num_slides-1}: Content slides with bullet points (max 4 bullets per slide)
{num_slides}: Conclusion slide with key takeaways

For each slide, provide:
- slide_title: The slide title
- bullet_points: List of 2-4 concise bullet points
- speaker_notes: Brief speaker notes for that slide

Format your response as a valid JSON object with this structure:
{{
    "presentation_title": "Title of the presentation",
    "presentation_subtitle": "Subtitle or tagline",
    "slides": [
        {{
            "slide_number": 1,
            "slide_title": "Title Slide",
            "bullet_points": ["Point 1", "Point 2"],
            "speaker_notes": "Brief notes"
        }}
    ]
}}

Make it informative, engaging, and professional."""

        response = requests.post(
            f"{OLLAMA_API_URL}/api/generate",
            json={
                "model": OLLAMA_MODEL,
                "prompt": prompt,
                "stream": False,
                "temperature": 0.7,
            },
            timeout=120
        )
        
        if response.status_code != 200:
            raise Exception(f"LLM API error: {response.status_code}")
        
        result = response.json()
        response_text = result.get("response", "")
        
        # Parse JSON from response
        try:
            # Try to extract JSON from the response
            json_start = response_text.find('{')
            json_end = response_text.rfind('}') + 1
            if json_start >= 0 and json_end > json_start:
                json_str = response_text[json_start:json_end]
                presentation_data = json.loads(json_str)
                return presentation_data
        except json.JSONDecodeError:
            pass
        
        # Fallback if JSON parsing fails
        return generate_fallback_presentation(topic, num_slides)
        
    except Exception as e:
        print(f"Error querying LLM: {str(e)}")
        return generate_fallback_presentation(topic, num_slides)


def generate_fallback_presentation(topic: str, num_slides: int) -> Dict[str, Any]:
    """
    Generate a fallback presentation structure if LLM fails
    """
    slides = []
    
    # Title slide
    slides.append({
        "slide_number": 1,
        "slide_title": "Presentation",
        "bullet_points": [topic, "A Comprehensive Overview"],
        "speaker_notes": f"Welcome to the presentation on {topic}"
    })
    
    # Content slides
    for i in range(2, num_slides):
        slides.append({
            "slide_number": i,
            "slide_title": f"{topic} - Part {i-1}",
            "bullet_points": [
                f"Key point {i-1}-1",
                f"Key point {i-1}-2",
                f"Key point {i-1}-3"
            ],
            "speaker_notes": f"Details about part {i-1} of {topic}"
        })
    
    # Conclusion slide
    slides.append({
        "slide_number": num_slides,
        "slide_title": "Thank You",
        "bullet_points": [
            "Summary of key points",
            "Thank you for your attention",
            "Questions?"
        ],
        "speaker_notes": "Closing remarks and Q&A session"
    })
    
    return {
        "presentation_title": topic,
        "presentation_subtitle": "A Comprehensive Overview",
        "slides": slides
    }


def create_presentation_file(presentation_data: Dict[str, Any], output_path: str) -> str:
    """
    Create a PPTX presentation file from structured data
    
    Args:
        presentation_data: Dictionary with presentation structure
        output_path: Path where to save the PPTX file
    
    Returns:
        Path to the created presentation file
    """
    try:
        # Create presentation object
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        title = presentation_data.get("presentation_title", "Presentation")
        subtitle = presentation_data.get("presentation_subtitle", "")
        slides_data = presentation_data.get("slides", [])
        
        for slide_data in slides_data:
            slide_number = slide_data.get("slide_number", 1)
            slide_title = slide_data.get("slide_title", "")
            bullet_points = slide_data.get("bullet_points", [])
            
            # Create slide layout
            if slide_number == 1:
                # Title slide
                slide_layout = prs.slide_layouts[6]  # Blank layout
                slide = prs.slides.add_slide(slide_layout)
                
                # Add background color
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = COLOR_PRIMARY
                
                # Add title
                title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
                title_frame = title_box.text_frame
                title_frame.word_wrap = True
                p = title_frame.paragraphs[0]
                p.text = title
                p.font.size = Pt(54)
                p.font.bold = True
                p.font.color.rgb = COLOR_WHITE
                p.alignment = PP_ALIGN.CENTER
                
                # Add subtitle
                subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
                subtitle_frame = subtitle_box.text_frame
                subtitle_frame.word_wrap = True
                p = subtitle_frame.paragraphs[0]
                p.text = subtitle
                p.font.size = Pt(28)
                p.font.color.rgb = COLOR_SECONDARY
                p.alignment = PP_ALIGN.CENTER
                
                # Add date at bottom
                date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.5))
                date_frame = date_box.text_frame
                p = date_frame.paragraphs[0]
                p.text = datetime.now().strftime("%B %d, %Y")
                p.font.size = Pt(14)
                p.font.color.rgb = COLOR_WHITE
                p.alignment = PP_ALIGN.CENTER
            
            elif slide_number == len(slides_data):
                # Conclusion/Thank you slide
                slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(slide_layout)
                
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = COLOR_SECONDARY
                
                # Add title
                title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
                title_frame = title_box.text_frame
                p = title_frame.paragraphs[0]
                p.text = slide_title
                p.font.size = Pt(44)
                p.font.bold = True
                p.font.color.rgb = COLOR_WHITE
                p.alignment = PP_ALIGN.CENTER
                
                # Add bullet points
                content_box = slide.shapes.add_textbox(Inches(1.5), Inches(3.8), Inches(7), Inches(3))
                text_frame = content_box.text_frame
                text_frame.word_wrap = True
                
                for idx, point in enumerate(bullet_points):
                    if idx == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    
                    p.text = point
                    p.font.size = Pt(24)
                    p.font.color.rgb = COLOR_WHITE
                    p.level = 0
                    p.alignment = PP_ALIGN.CENTER
                    p.space_before = Pt(12)
            
            else:
                # Content slides
                slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(slide_layout)
                
                # Add background color
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = COLOR_WHITE
                
                # Add header bar
                header_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
                header_shape.fill.solid()
                header_shape.fill.fore_color.rgb = COLOR_PRIMARY
                header_shape.line.color.rgb = COLOR_PRIMARY
                
                # Add slide title
                title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.7))
                title_frame = title_box.text_frame
                title_frame.word_wrap = True
                p = title_frame.paragraphs[0]
                p.text = slide_title
                p.font.size = Pt(40)
                p.font.bold = True
                p.font.color.rgb = COLOR_WHITE
                
                # Add content
                content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5.5))
                text_frame = content_box.text_frame
                text_frame.word_wrap = True
                
                for idx, point in enumerate(bullet_points):
                    if idx == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    
                    p.text = point
                    p.font.size = Pt(20)
                    p.font.color.rgb = COLOR_TEXT
                    p.level = 0
                    p.space_before = Pt(12)
                
                # Add slide number at bottom right
                slide_num_box = slide.shapes.add_textbox(Inches(8.5), Inches(7), Inches(1.5), Inches(0.5))
                slide_num_frame = slide_num_box.text_frame
                p = slide_num_frame.paragraphs[0]
                p.text = f"{slide_number - 1} of {len(slides_data) - 1}"
                p.font.size = Pt(10)
                p.font.color.rgb = RGBColor(153, 153, 153)
                p.alignment = PP_ALIGN.RIGHT
        
        # Save presentation
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        prs.save(output_path)
        return output_path
        
    except Exception as e:
        raise Exception(f"Error creating presentation file: {str(e)}")


def generate_presentation(
    topic: str,
    num_slides: int = 10,
    style: str = "professional",
    output_dir: str = None
) -> Dict[str, Any]:
    """
    Main function to generate a complete presentation
    
    Args:
        topic: Topic for the presentation
        num_slides: Number of slides (default: 10)
        style: Presentation style (professional, educational, creative)
        output_dir: Output directory for the presentation file
    
    Returns:
        Dictionary with file path and presentation metadata
    """
    try:
        # Default output directory
        if not output_dir:
            output_dir = os.path.join(os.path.dirname(__file__), '..', 'storage', 'presentations')
        
        # Generate presentation content using LLM
        print(f"Generating presentation content for topic: {topic}")
        presentation_data = query_llm_for_presentation(topic, num_slides, style)
        
        # Create PPTX file
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        safe_topic = topic.replace(" ", "_").replace("/", "_")[:30]
        filename = f"{safe_topic}_{timestamp}.pptx"
        output_path = os.path.join(output_dir, filename)
        
        print(f"Creating presentation file: {output_path}")
        file_path = create_presentation_file(presentation_data, output_path)
        
        return {
            "success": True,
            "file_path": file_path,
            "filename": filename,
            "topic": topic,
            "num_slides": len(presentation_data.get("slides", [])),
            "presentation_title": presentation_data.get("presentation_title", ""),
            "message": f"Presentation '{presentation_data.get('presentation_title', topic)}' generated successfully"
        }
        
    except Exception as e:
        error_msg = f"Error generating presentation: {str(e)}"
        print(error_msg)
        return {
            "success": False,
            "error": error_msg,
            "topic": topic
        }
